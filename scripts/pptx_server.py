"""
AURA PPTX Server — FastAPI + python-pptx + fpdf2
Generates professional PowerPoint presentations and structured PDF reports.

Deploy on Ubuntu server:
  pip install fastapi uvicorn python-pptx requests fpdf2
  uvicorn pptx_server:app --host 0.0.0.0 --port 8200

Endpoints:
  POST /generate         — Generate a PPTX from structured content
  POST /generate-report  — Generate a structured PDF report
  POST /convert-to-pdf   — Convert PPTX to PDF (LibreOffice)
  GET  /health           — Health check

Supported PPTX layouts:
  title, section, content, two_column, image_right, image_left,
  image_full, key_metrics, quote, table, timeline, closing

Supported report sections:
  heading, paragraph, bullets, numbered_list, table,
  key_metrics, quote, page_break
"""

import os
import base64
import io
import subprocess
import tempfile
import requests
from datetime import datetime
from typing import Optional

from fastapi import FastAPI, HTTPException, Header
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

app = FastAPI(title="AURA PPTX Server")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

API_KEY = os.environ.get("PPTX_API_KEY", "aura-pptx-secret-key")
TEMPLATE_DIR = os.path.join(os.path.dirname(__file__), "templates")
os.makedirs(TEMPLATE_DIR, exist_ok=True)


# ─── Models ──────────────────────────────────────────────────
class SlideInput(BaseModel):
    title: str = ""
    layout: Optional[str] = "content"
    content: Optional[str] = None
    bullets: Optional[list[str]] = None
    image_url: Optional[str] = None
    image_caption: Optional[str] = None
    columns: Optional[list[dict]] = None
    key_metrics: Optional[list[dict]] = None
    quote: Optional[str] = None
    quote_author: Optional[str] = None
    table_data: Optional[dict] = None


class PresentationRequest(BaseModel):
    title: str
    slides: list[SlideInput]
    theme: Optional[str] = "professional"


class PresentationResponse(BaseModel):
    success: bool
    base64_data: str
    file_name: str
    slides_count: int
    size_bytes: int


# ─── Theme colors ────────────────────────────────────────────
THEMES = {
    "professional": {
        "primary": RGBColor(0x1B, 0x3A, 0x5C),
        "primary_dark": RGBColor(0x0F, 0x24, 0x40),
        "accent": RGBColor(0x34, 0x98, 0xDB),
        "accent2": RGBColor(0x2E, 0xCC, 0x71),
        "text_dark": RGBColor(0x2C, 0x3E, 0x50),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_subtle": RGBColor(0x7F, 0x8C, 0x8D),
        "bg_light": RGBColor(0xF8, 0xF9, 0xFA),
        "card_bg": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "minimal": {
        "primary": RGBColor(0x2D, 0x34, 0x36),
        "primary_dark": RGBColor(0x1E, 0x25, 0x27),
        "accent": RGBColor(0x00, 0xB8, 0x94),
        "accent2": RGBColor(0x00, 0xCE, 0xC9),
        "text_dark": RGBColor(0x2D, 0x34, 0x36),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_subtle": RGBColor(0xB2, 0xBE, 0xC3),
        "bg_light": RGBColor(0xFF, 0xFF, 0xFF),
        "card_bg": RGBColor(0xF5, 0xF6, 0xFA),
    },
    "corporate": {
        "primary": RGBColor(0x1A, 0x23, 0x7E),
        "primary_dark": RGBColor(0x0D, 0x16, 0x42),
        "accent": RGBColor(0xFF, 0x6F, 0x00),
        "accent2": RGBColor(0xFF, 0xA0, 0x00),
        "text_dark": RGBColor(0x1A, 0x23, 0x7E),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_subtle": RGBColor(0x9F, 0xA8, 0xDA),
        "bg_light": RGBColor(0xF5, 0xF5, 0xF5),
        "card_bg": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "modern": {
        "primary": RGBColor(0x1A, 0x1A, 0x2E),
        "primary_dark": RGBColor(0x0F, 0x0F, 0x1A),
        "accent": RGBColor(0xE9, 0x45, 0x60),
        "accent2": RGBColor(0x53, 0x3C, 0xD0),
        "text_dark": RGBColor(0xE0, 0xE0, 0xE0),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_subtle": RGBColor(0x88, 0x88, 0x88),
        "bg_light": RGBColor(0x16, 0x21, 0x3E),
        "card_bg": RGBColor(0x1A, 0x1A, 0x2E),
    },
    "creative": {
        "primary": RGBColor(0x6C, 0x5C, 0xE7),
        "primary_dark": RGBColor(0x48, 0x34, 0xD4),
        "accent": RGBColor(0xFD, 0x79, 0xA8),
        "accent2": RGBColor(0xFD, 0xCB, 0x6E),
        "text_dark": RGBColor(0x2D, 0x34, 0x36),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_subtle": RGBColor(0xB2, 0xBE, 0xC3),
        "bg_light": RGBColor(0xFA, 0xFA, 0xFA),
        "card_bg": RGBColor(0xFF, 0xFF, 0xFF),
    },
}


# ─── Helpers ─────────────────────────────────────────────────
def lighten_color(color: RGBColor, factor: float) -> RGBColor:
    r = min(255, int(color[0] + (255 - color[0]) * factor))
    g = min(255, int(color[1] + (255 - color[1]) * factor))
    b = min(255, int(color[2] + (255 - color[2]) * factor))
    return RGBColor(r, g, b)


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, diameter, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(diameter), Inches(diameter)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_shadow(shape):
    """Add a subtle drop shadow to a shape."""
    spPr = shape._element.spPr
    effectLst = spPr.makeelement(qn("a:effectLst"), {})
    outerShdw = effectLst.makeelement(qn("a:outerShdw"), {
        "blurRad": "76200", "dist": "25400", "dir": "5400000",
        "algn": "bl", "rotWithShape": "0",
    })
    srgbClr = outerShdw.makeelement(qn("a:srgbClr"), {"val": "000000"})
    alpha = srgbClr.makeelement(qn("a:alpha"), {"val": "15000"})
    srgbClr.append(alpha)
    outerShdw.append(srgbClr)
    effectLst.append(outerShdw)
    spPr.append(effectLst)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=RGBColor(0, 0, 0), bold=False, italic=False,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri",
                 vertical_anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    # Vertical alignment
    txBody = txBox._element.txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        anchor_map = {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}
        bodyPr.set("anchor", anchor_map.get(vertical_anchor, "t"))
    return txBox


def download_image(url: str) -> io.BytesIO:
    """Download image from URL. Returns BytesIO or None on failure."""
    try:
        resp = requests.get(url, timeout=8, stream=True,
                            headers={"User-Agent": "AURA-PPTX/1.0"})
        resp.raise_for_status()
        content_length = int(resp.headers.get("content-length", 0))
        if content_length > 10 * 1024 * 1024:  # 10MB max
            return None
        data = resp.content
        if len(data) > 10 * 1024 * 1024:
            return None
        buf = io.BytesIO(data)
        buf.seek(0)
        return buf
    except Exception:
        return None


def make_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_slide_header(slide, title, theme):
    """Standard header bar for content slides."""
    add_rect(slide, 0, 0, 0.12, 7.5, theme["primary"])
    add_rect(slide, 0.12, 0, 13.21, 1.4, theme["primary"])
    add_rounded_rect(slide, 0.5, 0.35, 0.5, 0.5, theme["accent"])
    add_text_box(
        slide, 1.3, 0.2, 11.0, 1.0, title,
        font_size=28, color=theme["text_light"], bold=True,
        font_name="Calibri Light", vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    add_rect(slide, 0.5, 1.4, 12.33, 0.04, theme["accent"])


def build_footer(slide, theme, slide_num, total_slides):
    add_rect(slide, 0.5, 7.0, 12.33, 0.02, lighten_color(theme["primary"], 0.7))
    add_circle(slide, 12.2, 6.8, 0.55, theme["accent"])
    add_text_box(slide, 12.2, 6.8, 0.55, 0.55, str(slide_num),
                 font_size=11, color=theme["text_light"], bold=True,
                 alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, 11.3, 7.05, 0.9, 0.3, f"/ {total_slides}",
                 font_size=9, color=theme["text_subtle"], alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, 0.5, 7.05, 1.5, 0.3, "AURA",
                 font_size=9, color=theme["text_subtle"], bold=True)


# ─── Layout: Title ──────────────────────────────────────────
def build_title_slide(prs, slide_data, title, theme, **_):
    slide = make_blank_slide(prs)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["primary_dark"]

    add_rect(slide, 0, 5.0, 13.33, 2.5, theme["primary"])
    add_circle(slide, 9.0, 3.5, 5.5, lighten_color(theme["primary"], 0.15))
    add_circle(slide, -1.0, -1.0, 3.5, lighten_color(theme["primary_dark"], 0.1))
    add_rect(slide, 0, 0, 13.33, 0.08, theme["accent"])
    add_rect(slide, 4.5, 3.55, 4.33, 0.06, theme["accent"])

    add_text_box(slide, 1.0, 1.5, 11.33, 1.8, slide_data.title,
                 font_size=44, color=theme["text_light"], bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Calibri Light")

    subtitle = slide_data.content or title
    add_text_box(slide, 2.0, 3.8, 9.33, 0.7, subtitle,
                 font_size=18, color=lighten_color(theme["text_light"], -0.2),
                 italic=True, alignment=PP_ALIGN.CENTER)

    date_str = datetime.now().strftime("%d %B %Y")
    add_text_box(slide, 1.0, 6.2, 11.33, 0.5,
                 f"AURA Assistant  |  {date_str}",
                 font_size=11, color=lighten_color(theme["text_light"], -0.3),
                 alignment=PP_ALIGN.CENTER)


# ─── Layout: Section divider ────────────────────────────────
def build_section_slide(prs, slide_data, theme, slide_num, total_slides, **_):
    slide = make_blank_slide(prs)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["primary"]

    add_circle(slide, 10.0, -1.5, 5.0, lighten_color(theme["primary"], 0.1))
    add_circle(slide, -2.0, 4.0, 4.0, lighten_color(theme["primary"], 0.08))
    add_rect(slide, 0, 0, 13.33, 0.08, theme["accent"])

    # Section number badge
    add_circle(slide, 6.1, 1.8, 1.1, theme["accent"])
    add_text_box(slide, 6.1, 1.8, 1.1, 1.1, str(slide_num),
                 font_size=28, color=theme["text_light"], bold=True,
                 alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, 1.0, 3.2, 11.33, 1.5, slide_data.title,
                 font_size=36, color=theme["text_light"], bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Calibri Light")

    add_rect(slide, 5.5, 4.9, 2.33, 0.06, theme["accent"])

    if slide_data.content:
        add_text_box(slide, 2.0, 5.2, 9.33, 0.8, slide_data.content,
                     font_size=16, color=lighten_color(theme["text_light"], -0.2),
                     italic=True, alignment=PP_ALIGN.CENTER)

    build_footer(slide, theme, slide_num, total_slides)


# ─── Layout: Content (bullets/text) ─────────────────────────
def build_content_slide(prs, slide_data, theme, slide_num, total_slides, **_):
    slide = make_blank_slide(prs)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["bg_light"]

    add_slide_header(slide, slide_data.title, theme)

    card = add_rounded_rect(slide, 0.5, 1.8, 12.33, 4.8, theme["card_bg"])
    add_shadow(card)

    if slide_data.bullets and len(slide_data.bullets) > 0:
        _add_bullets(slide, slide_data.bullets, theme)
    elif slide_data.content:
        _add_text(slide, slide_data.content, theme)

    build_footer(slide, theme, slide_num, total_slides)


def _add_bullets(slide, bullets, theme, x_start=1.0, y_start=2.15,
                 max_y=6.1, width=10.9):
    for i, text in enumerate(bullets):
        y = y_start + i * 0.72
        if y > max_y:
            break
        add_circle(slide, x_start, y + 0.1, 0.22, theme["accent"])
        add_text_box(slide, x_start + 0.45, y, width, 0.55, text,
                     font_size=16, color=theme["text_dark"],
                     vertical_anchor=MSO_ANCHOR.MIDDLE)
        if i < len(bullets) - 1 and y + 0.72 <= max_y:
            add_rect(slide, x_start + 0.45, y + 0.62, width - 0.4, 0.008,
                     lighten_color(theme["text_dark"], 0.85))


def _add_text(slide, content, theme, x=1.0, y=2.15, w=11.33, h=4.2):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(17)
    p.font.color.rgb = theme["text_dark"]
    p.font.name = "Calibri"
    p.line_spacing = Pt(26)
    p.space_after = Pt(10)


# ─── Layout: Two columns ────────────────────────────────────
def build_two_column_slide(prs, slide_data, theme, slide_num, total_slides, **_):
    slide = make_blank_slide(prs)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["bg_light"]

    add_slide_header(slide, slide_data.title, theme)

    columns = slide_data.columns or []
    if len(columns) < 2:
        # fallback to content slide
        build_content_slide.__wrapped__(prs, slide_data, theme, slide_num, total_slides)
        return

    for col_idx, col in enumerate(columns[:2]):
        x = 0.5 + col_idx * 6.25
        card = add_rounded_rect(slide, x, 1.8, 5.9, 4.8, theme["card_bg"])
        add_shadow(card)

        col_title = col.get("title", "")
        accent_color = theme["accent"] if col_idx == 0 else theme.get("accent2", theme["accent"])
        add_rect(slide, x, 1.8, 5.9, 0.06, accent_color)

        if col_title:
            add_text_box(slide, x + 0.3, 2.0, 5.3, 0.6, col_title,
                         font_size=20, color=theme["text_dark"], bold=True)

        col_bullets = col.get("bullets", [])
        for i, bt in enumerate(col_bullets):
            y = 2.7 + i * 0.62
            if y > 6.1:
                break
            add_circle(slide, x + 0.3, y + 0.08, 0.18, accent_color)
            add_text_box(slide, x + 0.7, y, 4.9, 0.5, bt,
                         font_size=14, color=theme["text_dark"],
                         vertical_anchor=MSO_ANCHOR.MIDDLE)

    build_footer(slide, theme, slide_num, total_slides)


# ─── Layout: Image + Text ───────────────────────────────────
def build_image_slide(prs, slide_data, theme, slide_num, total_slides, direction="right", **_):
    slide = make_blank_slide(prs)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["bg_light"]

    add_slide_header(slide, slide_data.title, theme)

    image_stream = None
    if slide_data.image_url:
        image_stream = download_image(slide_data.image_url)

    if direction == "right":
        text_x, text_w = 0.5, 7.0
        img_x, img_w = 8.0, 4.83
    else:
        text_x, text_w = 5.5, 7.33
        img_x, img_w = 0.5, 4.83

    # Text side
    card = add_rounded_rect(slide, text_x, 1.8, text_w, 4.8, theme["card_bg"])
    add_shadow(card)

    if slide_data.bullets:
        bx = text_x + 0.5
        _add_bullets(slide, slide_data.bullets, theme,
                     x_start=bx, y_start=2.15, max_y=6.1, width=text_w - 1.0)
    elif slide_data.content:
        _add_text(slide, slide_data.content, theme,
                  x=text_x + 0.3, y=2.15, w=text_w - 0.6, h=4.2)

    # Image side
    if image_stream:
        try:
            slide.shapes.add_picture(
                image_stream,
                Inches(img_x), Inches(1.8), Inches(img_w), Inches(4.8)
            )
        except Exception:
            _add_image_placeholder(slide, img_x, 1.8, img_w, 4.8, theme)
    else:
        _add_image_placeholder(slide, img_x, 1.8, img_w, 4.8, theme)

    if slide_data.image_caption:
        add_text_box(slide, img_x, 6.7, img_w, 0.3, slide_data.image_caption,
                     font_size=10, color=theme["text_subtle"],
                     italic=True, alignment=PP_ALIGN.CENTER)

    build_footer(slide, theme, slide_num, total_slides)


def _add_image_placeholder(slide, x, y, w, h, theme):
    """Add a placeholder rectangle when image download fails."""
    card = add_rounded_rect(slide, x, y, w, h, lighten_color(theme["primary"], 0.85))
    add_text_box(slide, x, y + h/2 - 0.3, w, 0.6, "Image",
                 font_size=14, color=theme["text_subtle"],
                 alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)


# ─── Layout: Full image ─────────────────────────────────────
def build_image_full_slide(prs, slide_data, theme, slide_num, total_slides, **_):
    slide = make_blank_slide(prs)

    image_stream = None
    if slide_data.image_url:
        image_stream = download_image(slide_data.image_url)

    if image_stream:
        try:
            slide.shapes.add_picture(
                image_stream,
                Inches(0), Inches(0), Inches(13.333), Inches(7.5)
            )
        except Exception:
            bg = slide.background
            bg.fill.solid()
            bg.fill.fore_color.rgb = theme["primary_dark"]
    else:
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = theme["primary_dark"]

    # Dark overlay band at bottom
    overlay = add_rect(slide, 0, 4.5, 13.33, 3.0, RGBColor(0x00, 0x00, 0x00))
    # Make overlay semi-transparent via XML
    spPr = overlay._element.spPr
    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is None:
        fillElem = spPr.find(f".//{qn('a:solidFill')}")
        if fillElem is not None:
            solidFill = fillElem
    if solidFill is not None:
        clr = solidFill.find(qn("a:srgbClr"))
        if clr is not None:
            alpha_el = clr.makeelement(qn("a:alpha"), {"val": "60000"})
            clr.append(alpha_el)

    add_text_box(slide, 1.0, 4.8, 11.33, 1.2, slide_data.title,
                 font_size=32, color=theme["text_light"], bold=True,
                 font_name="Calibri Light")

    if slide_data.content:
        add_text_box(slide, 1.0, 6.0, 11.33, 0.8, slide_data.content,
                     font_size=16, color=lighten_color(theme["text_light"], -0.15))

    build_footer(slide, theme, slide_num, total_slides)


# ─── Layout: Key metrics ────────────────────────────────────
def build_key_metrics_slide(prs, slide_data, theme, slide_num, total_slides, **_):
    slide = make_blank_slide(prs)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["bg_light"]

    add_slide_header(slide, slide_data.title, theme)

    metrics = slide_data.key_metrics or []
    count = min(len(metrics), 4)
    if count == 0:
        build_content_slide(prs, slide_data, theme, slide_num, total_slides)
        return

    total_width = 12.33
    gap = 0.4
    card_w = (total_width - gap * (count - 1)) / count
    start_x = 0.5

    colors = [theme["accent"], theme.get("accent2", theme["accent"]),
              theme["primary"], lighten_color(theme["accent"], 0.3)]

    for i, metric in enumerate(metrics[:4]):
        x = start_x + i * (card_w + gap)
        card = add_rounded_rect(slide, x, 2.2, card_w, 3.8, theme["card_bg"])
        add_shadow(card)

        # Top color bar
        add_rect(slide, x, 2.2, card_w, 0.08, colors[i % len(colors)])

        value = metric.get("value", "0")
        label = metric.get("label", "")

        # Big value
        add_text_box(slide, x, 2.8, card_w, 1.6, value,
                     font_size=42, color=colors[i % len(colors)], bold=True,
                     alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
                     font_name="Calibri Light")

        # Separator
        add_rect(slide, x + card_w * 0.2, 4.5, card_w * 0.6, 0.03,
                 lighten_color(theme["text_dark"], 0.8))

        # Label
        add_text_box(slide, x + 0.2, 4.7, card_w - 0.4, 1.0, label,
                     font_size=14, color=theme["text_dark"],
                     alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.TOP)

    build_footer(slide, theme, slide_num, total_slides)


# ─── Layout: Quote ───────────────────────────────────────────
def build_quote_slide(prs, slide_data, theme, slide_num, total_slides, **_):
    slide = make_blank_slide(prs)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["bg_light"]

    card_y = 2.0
    if slide_data.title:
        add_slide_header(slide, slide_data.title, theme)
    else:
        card_y = 1.0

    # Quote card
    card = add_rounded_rect(slide, 1.0, card_y, 11.33, 4.5, theme["card_bg"])
    add_shadow(card)

    # Left accent bar on card
    add_rect(slide, 1.0, card_y, 0.12, 4.5, theme["accent"])

    # Big quote mark
    add_text_box(slide, 1.5, card_y, 1.5, 1.5, "\u201C",
                 font_size=72, color=lighten_color(theme["accent"], 0.5),
                 font_name="Georgia", bold=True)

    # Quote text
    quote_text = slide_data.quote or slide_data.content or ""
    add_text_box(slide, 2.0, card_y + 0.8, 9.5, 2.5, quote_text,
                 font_size=22, color=theme["text_dark"], italic=True,
                 font_name="Georgia")

    # Author
    author = slide_data.quote_author or ""
    if author:
        add_rect(slide, 2.0, card_y + 3.5, 1.5, 0.04, theme["accent"])
        add_text_box(slide, 2.0, card_y + 3.7, 9.5, 0.5, f"— {author}",
                     font_size=16, color=theme["accent"], bold=True)

    build_footer(slide, theme, slide_num, total_slides)


# ─── Layout: Table ───────────────────────────────────────────
def build_table_slide(prs, slide_data, theme, slide_num, total_slides, **_):
    slide = make_blank_slide(prs)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["bg_light"]

    add_slide_header(slide, slide_data.title, theme)

    table_data = slide_data.table_data or {}
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])

    if not headers:
        build_content_slide(prs, slide_data, theme, slide_num, total_slides)
        return

    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header

    table_width = Inches(12.33)
    table_height = Inches(min(4.5, 0.5 + num_rows * 0.5))
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.8), table_width, table_height
    )
    table = table_shape.table

    # Style header row
    for j, hdr in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme["primary"]
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = theme["text_light"]
        p.font.bold = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

    # Style data rows
    for i, row in enumerate(rows):
        row_color = theme["card_bg"] if i % 2 == 0 else lighten_color(theme["bg_light"], 0.3)
        for j, val in enumerate(row[:num_cols]):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_color
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.color.rgb = theme["text_dark"]
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER

    build_footer(slide, theme, slide_num, total_slides)


# ─── Layout: Timeline ───────────────────────────────────────
def build_timeline_slide(prs, slide_data, theme, slide_num, total_slides, **_):
    slide = make_blank_slide(prs)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["bg_light"]

    add_slide_header(slide, slide_data.title, theme)

    bullets = slide_data.bullets or []
    count = min(len(bullets), 5)
    if count == 0:
        build_content_slide(prs, slide_data, theme, slide_num, total_slides)
        return

    # Horizontal timeline line
    line_y = 3.5
    add_rect(slide, 0.8, line_y, 11.73, 0.06, theme["accent"])

    step_width = 11.73 / count
    colors = [theme["accent"], theme.get("accent2", theme["accent"]),
              theme["primary"], lighten_color(theme["accent"], 0.3),
              lighten_color(theme["primary"], 0.3)]

    for i in range(count):
        cx = 0.8 + i * step_width + step_width / 2

        # Circle on timeline
        circle_d = 0.7
        add_circle(slide, cx - circle_d/2, line_y - circle_d/2 + 0.03,
                   circle_d, colors[i % len(colors)])
        # Step number
        add_text_box(slide, cx - circle_d/2, line_y - circle_d/2 + 0.03,
                     circle_d, circle_d, str(i + 1),
                     font_size=18, color=theme["text_light"], bold=True,
                     alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Step text below
        add_text_box(slide, cx - step_width/2, 4.2, step_width, 2.3,
                     bullets[i],
                     font_size=13, color=theme["text_dark"],
                     alignment=PP_ALIGN.CENTER)

    build_footer(slide, theme, slide_num, total_slides)


# ─── Layout: Closing ────────────────────────────────────────
def build_closing_slide(prs, theme, **_):
    slide = make_blank_slide(prs)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["primary_dark"]

    add_circle(slide, -2.0, 2.0, 6.0, lighten_color(theme["primary_dark"], 0.08))
    add_circle(slide, 10.0, -1.0, 4.5, lighten_color(theme["primary_dark"], 0.06))
    add_rect(slide, 0, 0, 13.33, 0.08, theme["accent"])

    add_text_box(slide, 1.0, 2.5, 11.33, 1.5, "Merci",
                 font_size=48, color=theme["text_light"], bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Calibri Light")
    add_rect(slide, 5.5, 4.2, 2.33, 0.06, theme["accent"])
    add_text_box(slide, 1.0, 4.6, 11.33, 0.6, "Généré par AURA Assistant",
                 font_size=14, color=lighten_color(theme["text_light"], -0.3),
                 italic=True, alignment=PP_ALIGN.CENTER)


# ─── Layout router ──────────────────────────────────────────
LAYOUT_BUILDERS = {
    "title": build_title_slide,
    "section": build_section_slide,
    "content": build_content_slide,
    "two_column": build_two_column_slide,
    "image_right": build_image_slide,
    "image_left": build_image_slide,
    "image_full": build_image_full_slide,
    "key_metrics": build_key_metrics_slide,
    "quote": build_quote_slide,
    "table": build_table_slide,
    "timeline": build_timeline_slide,
    "closing": build_closing_slide,
}


# ─── Main generation ────────────────────────────────────────
def generate_presentation(req: PresentationRequest) -> tuple[bytes, int]:
    theme_name = req.theme or "professional"
    theme = THEMES.get(theme_name, THEMES["professional"])

    template_path = os.path.join(TEMPLATE_DIR, f"{theme_name}.pptx")
    if os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total_content_slides = max(len(req.slides) - 1, 1)

    for i, slide_data in enumerate(req.slides):
        layout = slide_data.layout or ("title" if i == 0 else "content")

        if layout == "title":
            build_title_slide(prs, slide_data, req.title, theme)
        elif layout in ("image_right", "image_left"):
            build_image_slide(prs, slide_data, theme, i, total_content_slides,
                              direction=layout.split("_")[1])
        elif layout == "closing":
            build_closing_slide(prs, theme)
        elif layout in LAYOUT_BUILDERS:
            LAYOUT_BUILDERS[layout](prs, slide_data, theme, i, total_content_slides)
        else:
            build_content_slide(prs, slide_data, theme, i, total_content_slides)

    build_closing_slide(prs, theme)

    total_slides = len(req.slides) + 1
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue(), total_slides


# ─── API endpoints ───────────────────────────────────────────
@app.get("/health")
async def health():
    return {"status": "ok", "service": "AURA PPTX Server", "version": "2.0"}


@app.post("/generate", response_model=PresentationResponse)
async def generate(req: PresentationRequest, x_api_key: str = Header(None)):
    if x_api_key != API_KEY:
        raise HTTPException(status_code=401, detail="Invalid API key")

    if not req.slides:
        raise HTTPException(status_code=400, detail="slides[] is required")

    try:
        pptx_bytes, slides_count = generate_presentation(req)
        b64 = base64.b64encode(pptx_bytes).decode("utf-8")
        safe_title = "".join(
            c for c in req.title if c.isalnum() or c in " -_"
        ).strip().replace(" ", "_")[:50]

        return PresentationResponse(
            success=True,
            base64_data=b64,
            file_name=f"{safe_title}.pptx",
            slides_count=slides_count,
            size_bytes=len(pptx_bytes),
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ─── PDF Conversion ──────────────────────────────────────────
class ConvertRequest(BaseModel):
    base64_data: str  # PPTX en base64


class ConvertResponse(BaseModel):
    success: bool
    base64_data: str  # PDF en base64
    size_bytes: int


@app.post("/convert-to-pdf", response_model=ConvertResponse)
async def convert_to_pdf(req: ConvertRequest, x_api_key: str = Header(None)):
    if x_api_key != API_KEY:
        raise HTTPException(status_code=401, detail="Invalid API key")

    try:
        pptx_bytes = base64.b64decode(req.base64_data)

        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = os.path.join(tmpdir, "presentation.pptx")
            with open(pptx_path, "wb") as f:
                f.write(pptx_bytes)

            result = subprocess.run(
                [
                    "libreoffice", "--headless", "--convert-to", "pdf",
                    "--outdir", tmpdir, pptx_path,
                ],
                check=True,
                timeout=60,
                capture_output=True,
            )
            print(f"[convert-to-pdf] LibreOffice stdout: {result.stdout.decode()}")

            pdf_path = os.path.join(tmpdir, "presentation.pdf")
            if not os.path.exists(pdf_path):
                raise FileNotFoundError(f"PDF not generated. stderr: {result.stderr.decode()}")

            with open(pdf_path, "rb") as f:
                pdf_bytes = f.read()

        return ConvertResponse(
            success=True,
            base64_data=base64.b64encode(pdf_bytes).decode("utf-8"),
            size_bytes=len(pdf_bytes),
        )
    except subprocess.TimeoutExpired:
        raise HTTPException(status_code=500, detail="PDF conversion timed out")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ─── Report PDF Generation (fpdf2) ──────────────────────────
from fpdf import FPDF
import colorsys
import tempfile
import urllib.request

FONT_DIR = os.path.join(os.path.dirname(__file__), "fonts")


# ─── Dynamic Color System ────────────────────────────────────

COLOR_NAMES = {
    # French
    "rouge": "#E74C3C", "bleu": "#3498DB", "vert": "#2ECC71",
    "orange": "#FF6F00", "violet": "#6C5CE7", "noir": "#2D3436",
    "rose": "#FD79A8", "jaune": "#F1C40F", "turquoise": "#00B894",
    "gris": "#636E72", "bordeaux": "#6B0848", "marine": "#1B3A5C",
    "corail": "#FF6B6B", "indigo": "#4834D4", "emeraude": "#00B16A",
    # English
    "red": "#E74C3C", "blue": "#3498DB", "green": "#2ECC71",
    "purple": "#6C5CE7", "black": "#2D3436", "pink": "#FD79A8",
    "yellow": "#F1C40F", "gray": "#636E72", "grey": "#636E72",
    "teal": "#00B894", "coral": "#FF6B6B", "navy": "#1B3A5C",
}


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    """Convert #RRGGBB or RRGGBB to (r, g, b) tuple."""
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = h[0]*2 + h[1]*2 + h[2]*2
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _darken(color: tuple, factor: float) -> tuple[int, int, int]:
    """Darken a color by factor (0-1)."""
    return (
        max(0, int(color[0] * (1 - factor))),
        max(0, int(color[1] * (1 - factor))),
        max(0, int(color[2] * (1 - factor))),
    )


def _hue_shift(color: tuple, degrees: float) -> tuple[int, int, int]:
    """Shift hue of a color by given degrees."""
    r, g, b = color[0] / 255.0, color[1] / 255.0, color[2] / 255.0
    h, s, v = colorsys.rgb_to_hsv(r, g, b)
    h = (h + degrees / 360.0) % 1.0
    r2, g2, b2 = colorsys.hsv_to_rgb(h, s, v)
    return (int(r2 * 255), int(g2 * 255), int(b2 * 255))


def generate_theme_from_color(color_input: str) -> dict:
    """Generate a full theme from a single color (hex or name)."""
    color_input = color_input.strip().lower()
    if color_input in COLOR_NAMES:
        color_input = COLOR_NAMES[color_input]
    if not color_input.startswith("#"):
        color_input = "#" + color_input

    primary = _hex_to_rgb(color_input)
    return {
        "primary": primary,
        "primary_dark": _darken(primary, 0.3),
        "accent": primary,
        "accent2": _hue_shift(primary, 30),
        "text_dark": (44, 62, 80),
        "text_light": (255, 255, 255),
        "text_subtle": (127, 140, 141),
        "bg_light": (248, 249, 250),
        "card_bg": (255, 255, 255),
    }


# ─── Document Types ──────────────────────────────────────────

DOCUMENT_TYPES = {
    "rapport_intervention": {
        "label": "Rapport d'intervention",
        "template": "executive",
        "default_footer": "Document interne — Rapport d'intervention",
        "required_metadata": ["Date d'intervention", "Intervenant", "Client/Site", "Objet"],
    },
    "brief_technique": {
        "label": "Brief technique",
        "template": "executive",
        "default_footer": "Document interne — Usage réservé aux destinataires mentionnés",
        "required_metadata": ["Destinataires", "Date", "Statut", "Durée estimée"],
    },
    "recap_brief": {
        "label": "Récapitulatif de brief",
        "template": "executive",
        "default_footer": "Récapitulatif — Ne pas diffuser",
        "required_metadata": ["Client", "Projet", "Date", "Participants"],
    },
    "convention_publicitaire": {
        "label": "Convention d'accompagnement publicitaire",
        "template": "executive",
        "default_footer": "Convention confidentielle",
        "required_metadata": ["Annonceur", "Agence", "Date de signature", "Durée"],
    },
    "analyse": {
        "label": "Analyse / Étude",
        "template": "executive",
        "default_footer": "Document d'analyse interne",
        "required_metadata": ["Auteur", "Date", "Département"],
    },
    "compte_rendu": {
        "label": "Compte-rendu de réunion",
        "template": "executive",
        "default_footer": "Compte-rendu — Diffusion interne",
        "required_metadata": ["Date", "Lieu", "Participants", "Animateur"],
    },
    "custom": {
        "label": "Document libre",
        "template": "executive",
        "default_footer": None,
        "required_metadata": [],
    },
}


# ─── Report Templates (layout, not colors) ──────────────────

REPORT_TEMPLATES = {
    "executive": {
        "cover_bg": "light",
        "cover_has_metadata": True,
        "heading_numbered": True,
        "heading_accent_bar": False,
        "paragraph_justified": True,
        "footer_text": "Document interne — Usage réservé aux destinataires mentionnés — Ne pas reproduire ni diffuser",
    },
    "modern": {
        "cover_bg": "dark",
        "cover_has_metadata": False,
        "heading_numbered": False,
        "heading_accent_bar": True,
        "paragraph_justified": False,
        "footer_text": None,
    },
    "creative": {
        "cover_bg": "gradient",
        "cover_has_metadata": True,
        "heading_numbered": True,
        "heading_accent_bar": True,
        "paragraph_justified": False,
        "footer_text": None,
    },
}


# ─── Pydantic Models ─────────────────────────────────────────

class ReportSectionInput(BaseModel):
    type: str  # heading, paragraph, bullets, numbered_list, table, key_metrics, quote, page_break, info_box, alert_box, metadata_table, separator
    level: Optional[int] = 1
    text: Optional[str] = None
    items: Optional[list[str]] = None
    headers: Optional[list[str]] = None
    rows: Optional[list[list[str]]] = None
    metrics: Optional[list[dict]] = None
    author: Optional[str] = None
    box_type: Optional[str] = None        # For alert_box: warning, info, tip, security, hardware, forbidden
    box_title: Optional[str] = None       # For info_box / alert_box
    metadata: Optional[list[dict]] = None # For metadata_table: [{key, value}]


class ReportRequest(BaseModel):
    title: str
    subtitle: Optional[str] = None
    theme: Optional[str] = "professional"
    template: Optional[str] = "executive"
    document_type: Optional[str] = "custom"
    custom_color: Optional[str] = None
    metadata: Optional[list[dict]] = None    # Cover page metadata [{key, value}]
    footer_text: Optional[str] = None
    logo_url: Optional[str] = None
    include_logo: Optional[bool] = True
    reference: Optional[str] = None  # Document reference (e.g. RI-2026-0323-NOM)
    sections: list[ReportSectionInput]


class ReportResponse(BaseModel):
    success: bool
    base64_data: str
    file_name: str
    pages_count: int
    size_bytes: int


# Convert RGBColor (pptx) to (r, g, b) tuple for fpdf2
def rgb_tuple(color: RGBColor) -> tuple[int, int, int]:
    return (color[0], color[1], color[2])


# Theme colors as plain tuples for PDF
def get_pdf_theme(theme_name: str) -> dict:
    src = THEMES.get(theme_name, THEMES["professional"])
    return {k: rgb_tuple(v) for k, v in src.items()}


class AuraReport(FPDF):
    """Custom PDF class with professional header/footer."""

    def __init__(self, title: str, subtitle: str, theme: dict):
        super().__init__(orientation="P", unit="mm", format="A4")
        self.report_title = title
        self.report_subtitle = subtitle
        self.theme = theme
        self.is_cover_page = True
        self.template_config: dict = REPORT_TEMPLATES["executive"]
        # Heading auto-numbering counters
        self.h1_counter = 0
        self.h2_counter = 0
        self.h3_counter = 0
        # Logo and reference for header reuse
        self.logo_path_local: Optional[str] = None
        self.reference: str = ""
        # Register Unicode fonts
        self.add_font("DejaVu", "", os.path.join(FONT_DIR, "DejaVuSans.ttf"))
        self.add_font("DejaVu", "B", os.path.join(FONT_DIR, "DejaVuSans-Bold.ttf"))
        self.add_font("DejaVu", "I", os.path.join(FONT_DIR, "DejaVuSans-Oblique.ttf"))
        self.add_font("DejaVu", "BI", os.path.join(FONT_DIR, "DejaVuSans-BoldOblique.ttf"))
        self.set_auto_page_break(auto=True, margin=25)
        self.alias_nb_pages()

    def header(self):
        if self.is_cover_page:
            return
        # Full primary dark header bar (15mm)
        self.set_fill_color(*self.theme["primary_dark"])
        self.rect(0, 0, 210, 15, "F")
        # Accent bar below
        self.set_fill_color(*self.theme["accent"])
        self.rect(0, 15, 210, 1.5, "F")
        # Logo in header (small)
        if self.logo_path_local:
            try:
                self.image(self.logo_path_local, x=8, y=1.5, w=15)
            except Exception:
                pass
        # Report title in white, right-aligned
        self.set_font("DejaVu", "B", 9)
        self.set_text_color(*self.theme["text_light"])
        self.set_xy(80, 3)
        self.cell(122, 5, self.report_title.upper(), align="R")
        # Reference below title
        if self.reference:
            self.set_font("DejaVu", "", 7)
            self.set_xy(80, 8.5)
            self.cell(122, 4, f"Ref. : {self.reference}", align="R")
        self.ln(20)

    def footer(self):
        if self.is_cover_page:
            return
        self.set_y(-14)
        # Accent bar
        self.set_fill_color(*self.theme["accent"])
        self.rect(0, self.get_y(), 210, 0.8, "F")
        # Primary footer bar
        self.set_fill_color(*self.theme["primary_dark"])
        self.rect(0, self.get_y() + 0.8, 210, 13.2, "F")
        # Footer text (left)
        custom_footer = self.template_config.get("footer_text")
        self.set_font("DejaVu", "", 7)
        self.set_text_color(*self.theme["text_light"])
        self.set_xy(10, self.get_y() + 3.5)
        if custom_footer:
            self.cell(140, 5, custom_footer, align="L")
        else:
            self.cell(140, 5, "AURA", align="L")
        # Page number (right)
        self.set_font("DejaVu", "", 7)
        self.set_xy(150, self.get_y())
        self.cell(50, 5, f"Page {self.page_no()} / {{nb}}", align="R")


def _lighten(color: tuple, factor: float) -> tuple[int, int, int]:
    """Lighten a color tuple by factor (0-1)."""
    return (
        min(255, int(color[0] + (255 - color[0]) * factor)),
        min(255, int(color[1] + (255 - color[1]) * factor)),
        min(255, int(color[2] + (255 - color[2]) * factor)),
    )


def _download_logo(url: str) -> Optional[str]:
    """Download logo from URL to a temp file. Returns path or None."""
    try:
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        urllib.request.urlretrieve(url, tmp.name)
        return tmp.name
    except Exception as e:
        print(f"[report] Logo download failed: {e}")
        return None


def _render_cover_header(pdf: AuraReport, req: ReportRequest, theme: dict, text_color: tuple):
    """Shared cover header: logo LEFT, doc type + reference RIGHT, accent bar below."""
    # Logo on LEFT (use pre-downloaded logo)
    if pdf.logo_path_local:
        try:
            pdf.image(pdf.logo_path_local, x=15, y=8, w=30)
        except Exception as e:
            print(f"[report] Logo insert failed: {e}")

    # Document type label on RIGHT
    doc_type = DOCUMENT_TYPES.get(req.document_type or "custom", DOCUMENT_TYPES["custom"])
    label = doc_type.get("label", "")
    if label:
        pdf.set_font("DejaVu", "B", 10)
        pdf.set_text_color(*text_color)
        pdf.set_xy(80, 15)
        pdf.cell(115, 7, label.upper(), align="R")

    # Reference on RIGHT (below label)
    if req.reference:
        pdf.set_font("DejaVu", "", 8)
        pdf.set_text_color(*text_color)
        pdf.set_xy(80, 23)
        pdf.cell(115, 6, f"Ref. : {req.reference}", align="R")

    # Accent bar below header
    pdf.set_fill_color(*theme["accent"])
    pdf.rect(0, 45, 210, 1.5, "F")


def _render_cover_executive(pdf: AuraReport, req: ReportRequest, theme: dict):
    """Clean white cover — like Claude artifacts."""
    # Top accent bar
    pdf.set_fill_color(*theme["primary"])
    pdf.rect(0, 0, 210, 3, "F")
    pdf.set_fill_color(*theme["accent"])
    pdf.rect(0, 3, 210, 0.8, "F")

    # Header: logo LEFT, doc type RIGHT
    _render_cover_header(pdf, req, theme, theme["text_subtle"])

    # Title
    pdf.set_font("DejaVu", "B", 20)
    pdf.set_text_color(*theme["primary"])
    pdf.set_xy(20, 52)
    pdf.multi_cell(165, 10, req.title, align="L")

    # Subtitle
    if req.subtitle:
        pdf.set_font("DejaVu", "I", 10)
        pdf.set_text_color(*theme["text_subtle"])
        pdf.set_xy(20, pdf.get_y() + 2)
        pdf.multi_cell(170, 6, req.subtitle, align="L")

    # Accent underline
    pdf.set_fill_color(*theme["accent"])
    pdf.rect(20, pdf.get_y() + 3, 40, 1.2, "F")

    # Metadata table on cover
    if req.metadata:
        y_meta = max(pdf.get_y() + 15, 130)
        pdf.set_xy(25, y_meta)
        key_w = 55
        val_w = 115
        row_h = 10
        pdf.set_draw_color(*_lighten(theme["text_subtle"], 0.4))
        pdf.set_line_width(0.3)
        for meta in req.metadata:
            y = pdf.get_y()
            if y > 260:
                break
            pdf.set_fill_color(*_lighten(theme["bg_light"], -0.05))
            pdf.set_font("DejaVu", "B", 10)
            pdf.set_text_color(*theme["text_dark"])
            pdf.set_xy(25, y)
            pdf.cell(key_w, row_h, str(meta.get("key", "")), border=1, fill=True)
            pdf.set_fill_color(255, 255, 255)
            pdf.set_font("DejaVu", "", 10)
            pdf.set_xy(25 + key_w, y)
            pdf.cell(val_w, row_h, str(meta.get("value", "")), border=1, fill=True)
            pdf.ln(row_h)

    # Date and branding at bottom
    pdf.set_font("DejaVu", "", 10)
    pdf.set_text_color(*theme["text_subtle"])
    date_str = datetime.now().strftime("%d %B %Y")
    pdf.set_xy(20, 265)
    pdf.cell(170, 8, f"AURA  |  {date_str}", align="L")


def _render_cover_dark(pdf: AuraReport, req: ReportRequest, theme: dict):
    """Dark cover — original modern style."""
    pdf.set_fill_color(*theme["primary_dark"])
    pdf.rect(0, 0, 210, 297, "F")

    # Decorative circles
    light = _lighten(theme["primary_dark"], 0.1)
    pdf.set_fill_color(*light)
    pdf.ellipse(140, -20, 120, 120, "F")
    light2 = _lighten(theme["primary_dark"], 0.07)
    pdf.set_fill_color(*light2)
    pdf.ellipse(-30, 180, 100, 100, "F")

    # Top accent bar
    pdf.set_fill_color(*theme["accent"])
    pdf.rect(0, 0, 210, 2, "F")

    # Header: logo LEFT, doc type RIGHT (white text)
    _render_cover_header(pdf, req, theme, theme["text_light"])

    # Title
    pdf.set_font("DejaVu", "B", 22)
    pdf.set_text_color(*theme["text_light"])
    pdf.set_xy(20, 55)
    pdf.multi_cell(170, 11, req.title, align="L")

    # Accent line
    pdf.set_fill_color(*theme["accent"])
    pdf.rect(20, pdf.get_y() + 3, 40, 1.2, "F")

    if req.subtitle:
        pdf.set_font("DejaVu", "I", 11)
        pdf.set_text_color(*_lighten(theme["text_light"], -0.2))
        pdf.set_xy(20, pdf.get_y() + 8)
        pdf.multi_cell(170, 7, req.subtitle, align="L")

    pdf.set_font("DejaVu", "", 10)
    pdf.set_text_color(*_lighten(theme["text_light"], -0.3))
    date_str = datetime.now().strftime("%d %B %Y")
    pdf.set_xy(20, 230)
    pdf.cell(170, 8, f"AURA  |  {date_str}", align="L")

    pdf.set_fill_color(*theme["primary"])
    pdf.rect(0, 260, 210, 37, "F")


def _render_cover_gradient(pdf: AuraReport, req: ReportRequest, theme: dict):
    """Colorful accent cover — creative style."""
    pdf.set_fill_color(*theme["accent"])
    pdf.rect(0, 0, 210, 297, "F")

    # Decorative shapes
    lighter = _lighten(theme["accent"], 0.15)
    pdf.set_fill_color(*lighter)
    pdf.ellipse(130, -30, 140, 140, "F")
    darker = _darken(theme["accent"], 0.15)
    pdf.set_fill_color(*darker)
    pdf.ellipse(-40, 200, 120, 120, "F")

    # Header: logo LEFT, doc type RIGHT (white text)
    _render_cover_header(pdf, req, theme, theme["text_light"])

    # Title
    pdf.set_font("DejaVu", "B", 22)
    pdf.set_text_color(*theme["text_light"])
    pdf.set_xy(20, 55)
    pdf.multi_cell(170, 11, req.title, align="L")

    pdf.set_fill_color(*theme["text_light"])
    pdf.rect(20, pdf.get_y() + 3, 40, 1.2, "F")

    if req.subtitle:
        pdf.set_font("DejaVu", "I", 11)
        pdf.set_text_color(*_lighten(theme["text_light"], -0.1))
        pdf.set_xy(20, pdf.get_y() + 8)
        pdf.multi_cell(170, 7, req.subtitle, align="L")

    # Metadata on creative cover
    if req.metadata:
        y_meta = max(pdf.get_y() + 15, 150)
        pdf.set_xy(25, y_meta)
        for meta in req.metadata:
            y = pdf.get_y()
            if y > 250:
                break
            pdf.set_font("DejaVu", "B", 10)
            pdf.set_text_color(*theme["text_light"])
            pdf.set_xy(25, y)
            pdf.cell(50, 8, str(meta.get("key", "")))
            pdf.set_font("DejaVu", "", 10)
            pdf.cell(120, 8, str(meta.get("value", "")))
            pdf.ln(8)

    pdf.set_font("DejaVu", "", 11)
    pdf.set_text_color(*theme["text_light"])
    date_str = datetime.now().strftime("%d %B %Y")
    pdf.set_xy(20, 265)
    pdf.cell(170, 8, f"AURA  |  {date_str}", align="L")


def generate_report_pdf(req: ReportRequest) -> tuple[bytes, int]:
    # Resolve theme: custom_color takes priority
    if req.custom_color:
        theme = generate_theme_from_color(req.custom_color)
    else:
        theme = get_pdf_theme(req.theme or "professional")

    # Resolve template from document_type or explicit template
    doc_type = DOCUMENT_TYPES.get(req.document_type or "custom", DOCUMENT_TYPES["custom"])
    template_name = req.template or doc_type.get("template", "executive")
    template_config = REPORT_TEMPLATES.get(template_name, REPORT_TEMPLATES["executive"]).copy()

    # Override footer from request or document type
    if req.footer_text:
        template_config["footer_text"] = req.footer_text
    elif doc_type.get("default_footer") and not template_config.get("footer_text"):
        template_config["footer_text"] = doc_type["default_footer"]

    pdf = AuraReport(req.title, req.subtitle or "", theme)
    pdf.template_config = template_config
    pdf.reference = req.reference or ""

    # Download logo once for reuse in cover + headers
    if req.logo_url and req.include_logo is not False:
        pdf.logo_path_local = _download_logo(req.logo_url)

    # ── Cover page ──
    pdf.add_page()
    pdf.is_cover_page = True

    cover_bg = template_config.get("cover_bg", "light")
    if cover_bg == "dark":
        _render_cover_dark(pdf, req, theme)
    elif cover_bg == "gradient":
        _render_cover_gradient(pdf, req, theme)
    else:
        _render_cover_executive(pdf, req, theme)

    # ── Content pages ──
    pdf.is_cover_page = False

    for section in req.sections:
        _render_section(pdf, section, theme)

    # Output
    pdf_bytes = pdf.output()
    pages = pdf.page_no()

    # Cleanup logo temp file
    if pdf.logo_path_local:
        try:
            os.unlink(pdf.logo_path_local)
        except OSError:
            pass

    return bytes(pdf_bytes), pages


def _render_section(pdf: AuraReport, section: ReportSectionInput, theme: dict):
    s_type = section.type

    if s_type == "page_break":
        pdf.add_page()
        return

    if s_type == "heading":
        _render_heading(pdf, section, theme)
    elif s_type == "paragraph":
        _render_paragraph(pdf, section, theme)
    elif s_type == "bullets":
        _render_bullets(pdf, section, theme)
    elif s_type == "numbered_list":
        _render_numbered_list(pdf, section, theme)
    elif s_type == "table":
        _render_table(pdf, section, theme)
    elif s_type == "key_metrics":
        _render_key_metrics(pdf, section, theme)
    elif s_type == "quote":
        _render_quote(pdf, section, theme)
    elif s_type == "info_box":
        _render_info_box(pdf, section, theme)
    elif s_type == "alert_box":
        _render_alert_box(pdf, section, theme)
    elif s_type == "metadata_table":
        _render_metadata_table(pdf, section, theme)
    elif s_type == "separator":
        _render_separator(pdf, section, theme)


def _render_heading(pdf: AuraReport, section: ReportSectionInput, theme: dict):
    level = section.level or 1
    text = section.text or ""
    tpl = pdf.template_config

    # Auto-numbering
    prefix = ""
    if tpl.get("heading_numbered"):
        if level == 1:
            pdf.h1_counter += 1
            pdf.h2_counter = 0
            pdf.h3_counter = 0
            prefix = f"{pdf.h1_counter}. "
        elif level == 2:
            pdf.h2_counter += 1
            pdf.h3_counter = 0
            prefix = f"{pdf.h1_counter}.{pdf.h2_counter} "
        else:
            pdf.h3_counter += 1
            prefix = f"{pdf.h1_counter}.{pdf.h2_counter}.{pdf.h3_counter} "

    if level == 1:
        pdf.ln(8)
        if tpl.get("heading_accent_bar"):
            # Accent bar before H1 (modern/creative style)
            pdf.set_fill_color(*theme["accent"])
            pdf.rect(15, pdf.get_y(), 4, 10, "F")
            pdf.set_font("DejaVu", "B", 20)
            pdf.set_text_color(*theme["primary"])
            pdf.set_x(22)
            pdf.multi_cell(170, 10, prefix + text)
        else:
            # Executive style: colored number, no bar
            pdf.set_font("DejaVu", "B", 22)
            pdf.set_text_color(*theme["accent"])
            pdf.set_x(15)
            if prefix:
                prefix_w = pdf.get_string_width(prefix)
                pdf.cell(prefix_w, 10, prefix)
                pdf.set_text_color(*theme["primary"])
                pdf.set_font("DejaVu", "B", 20)
                pdf.multi_cell(180 - prefix_w, 10, text)
            else:
                pdf.set_text_color(*theme["primary"])
                pdf.set_font("DejaVu", "B", 20)
                pdf.multi_cell(180, 10, text)
        # Thin line under H1
        pdf.set_draw_color(*_lighten(theme["primary"], 0.7))
        pdf.set_line_width(0.3)
        pdf.line(15, pdf.get_y() + 1, 195, pdf.get_y() + 1)
        pdf.ln(5)
    elif level == 2:
        pdf.ln(5)
        pdf.set_font("DejaVu", "B", 15)
        pdf.set_text_color(*theme["accent"])
        pdf.set_x(15)
        pdf.multi_cell(175, 8, prefix + text)
        pdf.ln(2)
    else:
        pdf.ln(3)
        pdf.set_font("DejaVu", "BI", 12)
        pdf.set_text_color(*theme["accent"])
        pdf.set_x(15)
        pdf.multi_cell(175, 7, prefix + text)
        pdf.ln(2)


def _render_paragraph(pdf: AuraReport, section: ReportSectionInput, theme: dict):
    text = section.text or ""
    tpl = pdf.template_config
    pdf.set_font("DejaVu", "", 11)
    pdf.set_text_color(*theme["text_dark"])
    pdf.set_x(15)
    align = "J" if tpl.get("paragraph_justified") else "L"
    pdf.multi_cell(180, 7, text, align=align)
    pdf.ln(3)


def _render_bullets(pdf: AuraReport, section: ReportSectionInput, theme: dict):
    items = section.items or []
    for item in items:
        y = pdf.get_y()
        if y > 265:
            pdf.add_page()
            y = pdf.get_y()
        # Bullet character instead of ellipse
        pdf.set_font("DejaVu", "B", 11)
        pdf.set_text_color(*theme["accent"])
        pdf.set_xy(18, y)
        pdf.cell(5, 6.5, "\u2022")
        # Text
        pdf.set_font("DejaVu", "", 11)
        pdf.set_text_color(*theme["text_dark"])
        pdf.set_x(25)
        pdf.multi_cell(165, 6.5, item)
        pdf.ln(1.5)
    pdf.ln(2)


def _render_numbered_list(pdf: AuraReport, section: ReportSectionInput, theme: dict):
    items = section.items or []
    for i, item in enumerate(items, 1):
        y = pdf.get_y()
        if y > 265:
            pdf.add_page()
            y = pdf.get_y()
        # Number badge
        pdf.set_fill_color(*theme["accent"])
        pdf.set_font("DejaVu", "B", 9)
        pdf.set_text_color(*theme["text_light"])
        pdf.set_xy(16, y)
        pdf.cell(8, 7, str(i), align="C", fill=True)
        # Text
        pdf.set_font("DejaVu", "", 11)
        pdf.set_text_color(*theme["text_dark"])
        pdf.set_xy(27, y)
        pdf.multi_cell(163, 6.5, item)
        pdf.ln(1.5)
    pdf.ln(2)


def _render_table(pdf: AuraReport, section: ReportSectionInput, theme: dict):
    headers = section.headers or []
    rows = section.rows or []
    if not headers:
        return

    num_cols = len(headers)
    available_w = 180  # mm

    # Proportional column widths based on content length
    max_lens = [len(h) for h in headers]
    for row in rows:
        for j, val in enumerate(row[:num_cols]):
            if j < len(max_lens):
                max_lens[j] = max(max_lens[j], len(str(val)))
    total = max(sum(max_lens), 1)
    col_widths = [(l / total) * available_w for l in max_lens]
    # Clamp: min 20mm, max 90mm
    for i in range(len(col_widths)):
        col_widths[i] = max(20, min(90, col_widths[i]))
    # Normalize to fill available width
    cw_total = sum(col_widths)
    col_widths = [(w / cw_total) * available_w for w in col_widths]

    pdf.ln(3)
    x_start = 15
    line_h = 6  # height per text line

    pdf.set_line_width(0.2)

    def _draw_header_row():
        """Draw the header row (also used after page breaks)."""
        pdf.set_fill_color(*theme["primary"])
        pdf.set_text_color(*theme["text_light"])
        pdf.set_font("DejaVu", "B", 9)
        x = x_start
        y_h = pdf.get_y()
        # Measure header heights
        h_heights = []
        for j, hdr in enumerate(headers):
            tw = pdf.get_string_width(str(hdr))
            nlines = max(1, int(tw / (col_widths[j] - 4)) + 1)
            h_heights.append(nlines * line_h)
        max_hh = max(h_heights) + 4  # padding
        for j, hdr in enumerate(headers):
            pdf.set_fill_color(*theme["primary"])
            pdf.set_draw_color(*_lighten(theme["primary"], 0.3))
            pdf.rect(x, y_h, col_widths[j], max_hh, "DF")
            pdf.set_text_color(*theme["text_light"])
            pdf.set_font("DejaVu", "B", 9)
            pdf.set_xy(x + 2, y_h + 2)
            pdf.multi_cell(col_widths[j] - 4, line_h, str(hdr), border=0, align="C")
            x += col_widths[j]
        pdf.set_y(y_h + max_hh)

    _draw_header_row()

    # Data rows
    pdf.set_font("DejaVu", "", 9)
    row_bg_light = _lighten(theme["bg_light"], 0.0)
    row_bg_alt = _lighten(theme["primary"], 0.9)

    for i, row_data in enumerate(rows):
        # Pass 1: measure max row height
        row_heights = []
        pdf.set_font("DejaVu", "", 9)
        for j, val in enumerate(row_data[:num_cols]):
            text = str(val)
            tw = pdf.get_string_width(text)
            cell_w = col_widths[j] - 4  # padding
            nlines = max(1, int(tw / max(cell_w, 1)) + 1)
            row_heights.append(nlines * line_h)
        max_h = max(row_heights) + 4  # padding

        # Page break with header repeat
        if pdf.get_y() + max_h > 265:
            pdf.add_page()
            _draw_header_row()

        bg = row_bg_light if i % 2 == 0 else row_bg_alt
        x = x_start
        y_row = pdf.get_y()

        # Pass 2: draw cells with uniform height
        for j, val in enumerate(row_data[:num_cols]):
            # Draw filled+bordered rectangle
            pdf.set_fill_color(*bg)
            pdf.set_draw_color(*_lighten(theme["text_subtle"], 0.5))
            pdf.rect(x, y_row, col_widths[j], max_h, "DF")
            # Draw wrapped text inside
            pdf.set_text_color(*theme["text_dark"])
            pdf.set_font("DejaVu", "", 9)
            pdf.set_xy(x + 2, y_row + 2)
            pdf.multi_cell(col_widths[j] - 4, line_h, str(val), border=0, align="C")
            x += col_widths[j]
        pdf.set_y(y_row + max_h)

    pdf.ln(3)


def _render_key_metrics(pdf: AuraReport, section: ReportSectionInput, theme: dict):
    metrics = section.metrics or []
    if not metrics:
        return

    count = min(len(metrics), 4)
    available_w = 180
    gap = 5
    card_w = (available_w - gap * (count - 1)) / count
    x_start = 15

    colors = [
        theme["accent"],
        theme.get("accent2", theme["accent"]),
        theme["primary"],
        _lighten(theme["accent"], 0.3),
    ]

    y_top = pdf.get_y() + 3
    if y_top > 240:
        pdf.add_page()
        y_top = pdf.get_y() + 3

    card_h = 35

    for i, metric in enumerate(metrics[:4]):
        x = x_start + i * (card_w + gap)
        color = colors[i % len(colors)]

        # Card background
        pdf.set_fill_color(*_lighten(theme["bg_light"], 0.0))
        pdf.set_draw_color(*_lighten(theme["primary"], 0.7))
        pdf.set_line_width(0.3)
        pdf.rect(x, y_top, card_w, card_h, "DF")

        # Top accent bar on card
        pdf.set_fill_color(*color)
        pdf.rect(x, y_top, card_w, 2, "F")

        # Big value
        value = metric.get("value", "0")
        pdf.set_font("DejaVu", "B", 22)
        pdf.set_text_color(*color)
        pdf.set_xy(x, y_top + 6)
        pdf.cell(card_w, 10, value, align="C")

        # Separator
        pdf.set_fill_color(*_lighten(theme["text_dark"], 0.8))
        sep_w = card_w * 0.5
        pdf.rect(x + (card_w - sep_w) / 2, y_top + 19, sep_w, 0.3, "F")

        # Label
        label = metric.get("label", "")
        pdf.set_font("DejaVu", "", 9)
        pdf.set_text_color(*theme["text_dark"])
        pdf.set_xy(x + 2, y_top + 22)
        pdf.cell(card_w - 4, 6, label, align="C")

    pdf.set_y(y_top + card_h + 5)


def _render_quote(pdf: AuraReport, section: ReportSectionInput, theme: dict):
    text = section.text or ""
    author = section.author or ""

    pdf.ln(3)
    y_start = pdf.get_y()
    if y_start > 250:
        pdf.add_page()
        y_start = pdf.get_y()

    # Card background
    card_x = 20
    card_w = 170

    # Render text first to calculate height
    pdf.set_font("DejaVu", "I", 12)
    pdf.set_xy(card_x + 8, y_start + 5)
    # Save position to calculate actual height
    y_before = pdf.get_y()
    pdf.multi_cell(card_w - 16, 7, f"\u201c {text} \u201d")
    y_after_text = pdf.get_y()

    if author:
        pdf.ln(2)
        pdf.set_font("DejaVu", "B", 10)
        pdf.set_text_color(*theme["accent"])
        pdf.set_x(card_x + 8)
        pdf.cell(card_w - 16, 6, f"\u2014 {author}")
        y_after_text = pdf.get_y() + 8

    card_h = y_after_text - y_start + 5

    # Draw card background (behind text — we redraw)
    # fpdf2 doesn't support z-order, so we draw the accent bar at the position
    pdf.set_fill_color(*_lighten(theme["accent"], 0.85))
    pdf.rect(card_x, y_start, card_w, card_h, "F")

    # Left accent bar
    pdf.set_fill_color(*theme["accent"])
    pdf.rect(card_x, y_start, 3, card_h, "F")

    # Re-render text on top of background
    pdf.set_font("DejaVu", "I", 12)
    pdf.set_text_color(*theme["text_dark"])
    pdf.set_xy(card_x + 8, y_start + 5)
    pdf.multi_cell(card_w - 16, 7, f"\u201c {text} \u201d")

    if author:
        pdf.ln(2)
        # Accent line
        pdf.set_fill_color(*theme["accent"])
        pdf.rect(card_x + 8, pdf.get_y(), 20, 0.5, "F")
        pdf.ln(3)
        pdf.set_font("DejaVu", "B", 10)
        pdf.set_text_color(*theme["accent"])
        pdf.set_x(card_x + 8)
        pdf.cell(card_w - 16, 6, f"\u2014 {author}")

    pdf.ln(8)


# ─── New section renderers ───────────────────────────────────

ALERT_ICONS = {
    "warning": "\u26A0",   # ⚠
    "info": "\u2139",      # ℹ
    "tip": "\u2736",       # ✶
    "security": "\u25B2",  # ▲
    "hardware": "\u25A0",  # ■
    "forbidden": "\u2298", # ⊘
}

ALERT_COLORS = {
    "warning": (255, 193, 7),
    "info": (52, 152, 219),
    "tip": (46, 204, 113),
    "security": (231, 76, 60),
    "hardware": (108, 92, 231),
    "forbidden": (231, 76, 60),
}


def _render_info_box(pdf: AuraReport, section: ReportSectionInput, theme: dict):
    """Colored info box with optional title and content."""
    box_title = section.box_title or ""
    text = section.text or ""
    items = section.items or []

    pdf.ln(3)
    y_start = pdf.get_y()
    if y_start > 250:
        pdf.add_page()
        y_start = pdf.get_y()

    card_x = 15
    card_w = 180

    # First pass: measure height
    y_cursor = y_start + 5
    if box_title:
        pdf.set_font("DejaVu", "B", 12)
        pdf.set_xy(card_x + 8, y_cursor)
        pdf.multi_cell(card_w - 16, 7, box_title)
        y_cursor = pdf.get_y() + 2

    if text:
        pdf.set_font("DejaVu", "", 11)
        pdf.set_xy(card_x + 8, y_cursor)
        pdf.multi_cell(card_w - 16, 6.5, text)
        y_cursor = pdf.get_y()

    for item in items:
        pdf.set_font("DejaVu", "", 11)
        pdf.set_xy(card_x + 14, y_cursor)
        pdf.multi_cell(card_w - 22, 6.5, item)
        y_cursor = pdf.get_y() + 1

    card_h = y_cursor - y_start + 5

    # Draw background
    pdf.set_fill_color(*_lighten(theme["accent"], 0.85))
    pdf.rect(card_x, y_start, card_w, card_h, "F")
    # Left accent bar
    pdf.set_fill_color(*theme["accent"])
    pdf.rect(card_x, y_start, 3, card_h, "F")

    # Second pass: render text on top
    y_cursor = y_start + 5
    if box_title:
        pdf.set_font("DejaVu", "B", 12)
        pdf.set_text_color(*theme["primary"])
        pdf.set_xy(card_x + 8, y_cursor)
        pdf.multi_cell(card_w - 16, 7, box_title)
        y_cursor = pdf.get_y() + 2

    if text:
        pdf.set_font("DejaVu", "", 11)
        pdf.set_text_color(*theme["text_dark"])
        pdf.set_xy(card_x + 8, y_cursor)
        pdf.multi_cell(card_w - 16, 6.5, text)
        y_cursor = pdf.get_y()

    for item in items:
        pdf.set_fill_color(*theme["accent"])
        pdf.set_font("DejaVu", "B", 11)
        pdf.set_text_color(*theme["accent"])
        pdf.set_xy(card_x + 8, y_cursor + 0.5)
        pdf.cell(5, 6, "\u2022")
        pdf.set_font("DejaVu", "", 11)
        pdf.set_text_color(*theme["text_dark"])
        pdf.set_xy(card_x + 14, y_cursor)
        pdf.multi_cell(card_w - 22, 6.5, item)
        y_cursor = pdf.get_y() + 1

    pdf.set_y(y_start + card_h + 3)


def _render_alert_box(pdf: AuraReport, section: ReportSectionInput, theme: dict):
    """Alert box with icon and type-based coloring."""
    box_type = section.box_type or "info"
    text = section.text or ""
    box_title = section.box_title or ""
    icon = ALERT_ICONS.get(box_type, "\u2139")
    alert_color = ALERT_COLORS.get(box_type, (52, 152, 219))

    pdf.ln(2)
    y_start = pdf.get_y()
    if y_start > 255:
        pdf.add_page()
        y_start = pdf.get_y()

    card_x = 15
    card_w = 180
    icon_w = 12

    # Measure text height
    content = f"{box_title}\n{text}".strip() if box_title else text
    pdf.set_font("DejaVu", "", 10)
    pdf.set_xy(card_x + icon_w + 6, y_start + 4)
    pdf.multi_cell(card_w - icon_w - 14, 6, content)
    y_after = pdf.get_y()
    card_h = max(y_after - y_start + 4, 16)

    # Draw background
    pdf.set_fill_color(*_lighten(alert_color, 0.85))
    pdf.rect(card_x, y_start, card_w, card_h, "F")
    # Left border
    pdf.set_fill_color(*alert_color)
    pdf.rect(card_x, y_start, 3, card_h, "F")

    # Icon
    pdf.set_font("DejaVu", "B", 14)
    pdf.set_text_color(*alert_color)
    pdf.set_xy(card_x + 5, y_start + 3)
    pdf.cell(icon_w, 8, icon, align="C")

    # Re-render text
    y_cursor = y_start + 4
    if box_title:
        pdf.set_font("DejaVu", "B", 10)
        pdf.set_text_color(*_darken(alert_color, 0.2))
        pdf.set_xy(card_x + icon_w + 6, y_cursor)
        pdf.multi_cell(card_w - icon_w - 14, 6, box_title)
        y_cursor = pdf.get_y()

    if text:
        pdf.set_font("DejaVu", "", 10)
        pdf.set_text_color(*theme["text_dark"])
        pdf.set_xy(card_x + icon_w + 6, y_cursor)
        pdf.multi_cell(card_w - icon_w - 14, 6, text)

    pdf.set_y(y_start + card_h + 2)


def _render_metadata_table(pdf: AuraReport, section: ReportSectionInput, theme: dict):
    """Key-value metadata table."""
    metadata = section.metadata or []
    if not metadata:
        return

    pdf.ln(3)
    key_w = 55
    val_w = 125
    row_h = 9
    x_start = 15

    pdf.set_draw_color(*_lighten(theme["text_subtle"], 0.4))
    pdf.set_line_width(0.3)

    for meta in metadata:
        y = pdf.get_y()
        if y > 265:
            pdf.add_page()
        # Key cell
        pdf.set_fill_color(*_lighten(theme["bg_light"], -0.05))
        pdf.set_font("DejaVu", "B", 10)
        pdf.set_text_color(*theme["text_dark"])
        pdf.set_xy(x_start, pdf.get_y())
        pdf.cell(key_w, row_h, str(meta.get("key", "")), border=1, fill=True)
        # Value cell
        pdf.set_fill_color(255, 255, 255)
        pdf.set_font("DejaVu", "", 10)
        pdf.set_xy(x_start + key_w, pdf.get_y())
        pdf.cell(val_w, row_h, str(meta.get("value", "")), border=1, fill=True)
        pdf.ln(row_h)

    pdf.ln(3)


def _render_separator(pdf: AuraReport, section: ReportSectionInput, theme: dict):
    """Decorative line separator."""
    pdf.ln(5)
    y = pdf.get_y()
    pdf.set_draw_color(*_lighten(theme["primary"], 0.6))
    pdf.set_line_width(0.5)
    pdf.line(40, y, 170, y)
    pdf.ln(5)


# ─── Report API endpoint ────────────────────────────────────
@app.post("/generate-report", response_model=ReportResponse)
async def generate_report(req: ReportRequest, x_api_key: str = Header(None)):
    if x_api_key != API_KEY:
        raise HTTPException(status_code=401, detail="Invalid API key")

    if not req.sections:
        raise HTTPException(status_code=400, detail="sections[] is required")

    try:
        pdf_bytes, pages_count = generate_report_pdf(req)
        b64 = base64.b64encode(pdf_bytes).decode("utf-8")
        safe_title = "".join(
            c for c in req.title if c.isalnum() or c in " -_"
        ).strip().replace(" ", "_")[:50]

        return ReportResponse(
            success=True,
            base64_data=b64,
            file_name=f"{safe_title}.pdf",
            pages_count=pages_count,
            size_bytes=len(pdf_bytes),
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8200)
